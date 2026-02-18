"""
SVG 转 PPT 转换器
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.shapes.freeform import FreeformBuilder
from pptx.oxml import parse_xml
from lxml import etree
import re
import base64
from pathlib import Path


def px2emu(x):
    """SVG px 转 PPT EMU"""
    return int(x / 96 * 914400)


def hex2rgb(hex_color):
    """十六进制颜色转RGB"""
    if not hex_color or hex_color.lower() == "none":
        return None
    
    # 处理常见的颜色名称
    color_names = {
        'white': '#FFFFFF', 'black': '#000000', 'red': '#FF0000',
        'green': '#008000', 'blue': '#0000FF', 'yellow': '#FFFF00',
        'cyan': '#00FFFF', 'magenta': '#FF00FF', 'gray': '#808080',
        'grey': '#808080', 'orange': '#FFA500', 'purple': '#800080',
        'pink': '#FFC0CB', 'brown': '#A52A2A', 'lime': '#00FF00'
    }
    
    if hex_color.lower() in color_names:
        hex_color = color_names[hex_color.lower()]
    
    hex_color = hex_color.lstrip("#")
    
    # 检查是否是有效的十六进制颜色
    if not all(c in '0123456789ABCDEFabcdef' for c in hex_color):
        return None
    
    if len(hex_color) == 3:
        hex_color = "".join([c*2 for c in hex_color])
    
    if len(hex_color) != 6:
        return None
    
    return RGBColor(int(hex_color[0:2],16), int(hex_color[2:4],16), int(hex_color[4:6],16))


def add_arrow_to_line(connector, arrow_type='triangle'):
    """通过XML直接添加箭头到连接线"""
    try:
        spPr = None
        for element in connector._element.getchildren():
            if element.tag == "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr":
                spPr = element
                break
        
        if spPr is not None:
            ln = spPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is None:
                xml = f'<a:ln xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:tailEnd type="{arrow_type}"/></a:ln>'
                parsed_xml = parse_xml(xml)
                spPr.append(parsed_xml)
            else:
                xml = f'<a:tailEnd type="{arrow_type}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
                parsed_xml = parse_xml(xml)
                ln.append(parsed_xml)
    except Exception as e:
        print(f"添加箭头错误: {e}")


def parse_path_points(d):
    """解析 path 简单 M,L 直线命令"""
    points = []
    cmds = re.findall(r'([MLZ])([^MLZ]*)', d, re.IGNORECASE)
    current = (0,0)
    for cmd, vals in cmds:
        vals = [float(v) for v in re.findall(r'[-+]?\d*\.?\d+', vals)]
        if cmd.upper() == 'M' and len(vals) >= 2:
            current = (vals[0], vals[1])
            points.append(current)
        elif cmd.upper() == 'L':
            for i in range(0,len(vals),2):
                if i+1 < len(vals):
                    current = (vals[i], vals[i+1])
                    points.append(current)
        elif cmd.upper() == 'Z':
            pass
    return points


def svg_to_ppt(svg_path: str, output_path: str) -> str:
    """
    将SVG文件转换为PPT
    
    Args:
        svg_path: SVG文件路径
        output_path: 输出PPT文件路径
        
    Returns:
        输出PPT文件的路径
    """
    svg_path = Path(svg_path)
    output_path = Path(output_path)
    
    if not svg_path.exists():
        raise FileNotFoundError(f"SVG文件不存在: {svg_path}")
    
    # 确保输出目录存在
    output_path.parent.mkdir(parents=True, exist_ok=True)
    
    # 读取SVG
    tree = etree.parse(str(svg_path))
    root = tree.getroot()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for elem in root.iter():
        if not hasattr(elem, "tag") or elem.tag is None:
            continue
        if not isinstance(elem.tag, str):
            continue
        tag = etree.QName(elem.tag).localname

        # 跳过 defs 和 marker 内的元素
        parent = elem.getparent()
        if parent is not None:
            parent_tag = etree.QName(parent.tag).localname if isinstance(parent.tag, str) else None
            if parent_tag in ("defs", "marker"):
                continue

        # 图片处理
        if tag == "image":
            try:
                x = float(elem.get("x", 0))
                y = float(elem.get("y", 0))
                w = float(elem.get("width", 100))
                h = float(elem.get("height", 100))
                href = elem.get("{http://www.w3.org/1999/xlink}href") or elem.get("href")
                
                if href and href.startswith("data:image"):
                    match = re.match(r'data:image/(\w+);base64,(.+)', href)
                    if match:
                        img_format = match.group(1)
                        img_data = base64.b64decode(match.group(2))
                        
                        temp_img_path = f"/tmp/temp_img_{id(elem)}.{img_format}"
                        with open(temp_img_path, 'wb') as f:
                            f.write(img_data)
                        
                        slide.shapes.add_picture(temp_img_path, px2emu(x), px2emu(y), 
                                                width=px2emu(w), height=px2emu(h))
            except Exception as e:
                print(f"图片处理错误: {e}")
                continue

        # 矩形
        elif tag == "rect":
            x = float(elem.get("x",0))
            y = float(elem.get("y",0))
            w = float(elem.get("width",0))
            h = float(elem.get("height",0))
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px2emu(x), px2emu(y), px2emu(w), px2emu(h))
            fill = hex2rgb(elem.get("fill"))
            if fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill
            else:
                shape.fill.background()
            stroke = hex2rgb(elem.get("stroke"))
            if stroke:
                shape.line.color.rgb = stroke

        # 圆 / 椭圆
        elif tag in ("circle","ellipse"):
            cx = float(elem.get("cx",0))
            cy = float(elem.get("cy",0))
            if tag=="circle":
                r = float(elem.get("r",0))
                left = cx-r
                top = cy-r
                width = height = 2*r
            else:
                rx = float(elem.get("rx",0))
                ry = float(elem.get("ry",0))
                left = cx-rx
                top = cy-ry
                width = 2*rx
                height = 2*ry
            shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, px2emu(left), px2emu(top), px2emu(width), px2emu(height))
            fill = hex2rgb(elem.get("fill"))
            if fill:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill
            stroke = hex2rgb(elem.get("stroke"))
            if stroke:
                shape.line.color.rgb = stroke

        # 直线
        elif tag == "line":
            x1 = float(elem.get("x1",0))
            y1 = float(elem.get("y1",0))
            x2 = float(elem.get("x2",0))
            y2 = float(elem.get("y2",0))
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px2emu(x1), px2emu(y1), px2emu(x2), px2emu(y2))
            stroke = hex2rgb(elem.get("stroke"))
            if stroke:
                line.line.color.rgb = stroke
            marker_end = elem.get("marker-end", "")
            if "arrowhead" in marker_end:
                add_arrow_to_line(line, 'triangle')

        # polyline / polygon
        elif tag in ("polyline","polygon"):
            points = [float(v) for v in re.findall(r'[-+]?\d*\.?\d+', elem.get("points",""))]
            if len(points) >= 4:
                ff = FreeformBuilder(slide.shapes, px2emu(points[0]), px2emu(points[1]), 1.0, 1.0)
                for i in range(2,len(points),2):
                    if i+1 < len(points):
                        if i == len(points) - 2 and tag=="polygon":
                            ff.add_line_segments([(px2emu(points[i]), px2emu(points[i+1]))], close=True)
                        else:
                            ff.add_line_segments([(px2emu(points[i]), px2emu(points[i+1]))], close=False)
                shape = ff.convert_to_shape()
                fill = hex2rgb(elem.get("fill"))
                if fill:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = fill
                stroke = hex2rgb(elem.get("stroke"))
                if stroke:
                    shape.line.color.rgb = stroke

        # path
        elif tag == "path":
            d = elem.get("d","")
            pts = parse_path_points(d)
            if len(pts) >= 2:
                marker_end = elem.get("marker-end", "")
                has_arrow = "arrowhead" in marker_end
                
                for i in range(len(pts) - 1):
                    line = slide.shapes.add_connector(
                        MSO_CONNECTOR.STRAIGHT, 
                        px2emu(pts[i][0]), px2emu(pts[i][1]),
                        px2emu(pts[i+1][0]), px2emu(pts[i+1][1])
                    )
                    
                    stroke = hex2rgb(elem.get("stroke"))
                    if stroke:
                        line.line.color.rgb = stroke
                    
                    stroke_width = elem.get("stroke-width", "2")
                    try:
                        width_pt = float(stroke_width)
                        line.line.width = Pt(width_pt)
                    except:
                        pass
                    
                    if has_arrow and i == len(pts) - 2:
                        add_arrow_to_line(line, 'triangle')

        # text
        elif tag == "text":
            x = float(elem.get("x",0))
            y = float(elem.get("y",0))
            
            text_content = elem.text or ""
            for child in elem:
                if child.text:
                    text_content += child.text
                if child.tail:
                    text_content += child.tail
            
            if not text_content.strip():
                continue
                
            text_anchor = elem.get("text-anchor", "start")
            font_size = elem.get("font-size", "12")
            try:
                font_size_pt = int(re.findall(r'\d+', str(font_size))[0])
            except:
                font_size_pt = 12
            
            char_width = font_size_pt * 0.55
            estimated_width = len(text_content) * char_width
            
            text_x = x
            if text_anchor == "middle":
                text_x = x - estimated_width / 2
            elif text_anchor == "end":
                text_x = x - estimated_width
            
            text_y = y - font_size_pt * 0.75
            
            textbox = slide.shapes.add_textbox(
                px2emu(text_x), 
                px2emu(text_y), 
                px2emu(estimated_width + 20),
                px2emu(font_size_pt * 1.5)
            )
            tf = textbox.text_frame
            tf.text = text_content
            tf.word_wrap = False
            tf.paragraphs[0].font.size = Pt(font_size_pt)
            
            if text_anchor == "middle":
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            elif text_anchor == "end":
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
            else:
                tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            
            font_weight = elem.get("font-weight", "normal")
            if font_weight == "bold":
                tf.paragraphs[0].font.bold = True
            
            fill = hex2rgb(elem.get("fill"))
            if fill:
                tf.paragraphs[0].font.color.rgb = fill

    prs.save(str(output_path))
    print(f"SVG转PPT完成: {output_path}")
    return str(output_path)

