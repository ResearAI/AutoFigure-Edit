from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# from pptx.enum.dml import MSO_LINE_ARROWHEAD_STYLE
# from pptx.enum.dml import MSO_LINE_ARROWHEAD_STYLE
from lxml import etree
# 对应旧版的 MSO_ARROWHEAD_STYLE
class MSO_LINE_ARROWHEAD_STYLE:
    NONE = 0
    TRIANGLE = 1
    STEALTH = 2
    DIAMOND = 3
    OVAL = 4
    OPEN = 5

# 对应旧版的 MSO_ARROWHEAD_WIDTH
class MSO_LINE_ARROWHEAD_WIDTH:
    NARROW = 1
    MEDIUM = 2
    WIDE = 3

# 对应旧版的 MSO_ARROWHEAD_LENGTH
class MSO_LINE_ARROWHEAD_LENGTH:
    SHORT = 1
    MEDIUM = 2
    LONG = 3
EMU_PER_PX = 914400 / 96  # 96 DPI

def px_to_emu(px):
    return int(float(px) * EMU_PER_PX)

def hex_to_rgb(hex_color):
    hex_color = hex_color.replace("#", "")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )

def convert_svg_to_ppt(svg_string, output="output.pptx"):
    root = etree.fromstring(svg_string.encode())

    width = float(root.attrib.get("width", 1024))
    height = float(root.attrib.get("height", 1024))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    for elem in root.iter():
        tag = etree.QName(elem).localname

        # RECT
        if tag == "rect":
            x = float(elem.attrib.get("x", 0))
            y = float(elem.attrib.get("y", 0))
            w = float(elem.attrib.get("width"))
            h = float(elem.attrib.get("height"))

            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                px_to_emu(x),
                px_to_emu(y),
                px_to_emu(w),
                px_to_emu(h),
            )

            if "fill" in elem.attrib:
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(elem.attrib["fill"])

            if "stroke" in elem.attrib:
                shape.line.color.rgb = hex_to_rgb(elem.attrib["stroke"])

        # CIRCLE
        elif tag == "circle":
            cx = float(elem.attrib.get("cx"))
            cy = float(elem.attrib.get("cy"))
            r = float(elem.attrib.get("r"))

            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                px_to_emu(cx - r),
                px_to_emu(cy - r),
                px_to_emu(2 * r),
                px_to_emu(2 * r),
            )

            if "fill" in elem.attrib:
                shape.fill.solid()
                shape.fill.fore_color.rgb = hex_to_rgb(elem.attrib["fill"])

        # LINE
        elif tag == "line":
            x1 = float(elem.attrib.get("x1"))
            y1 = float(elem.attrib.get("y1"))
            x2 = float(elem.attrib.get("x2"))
            y2 = float(elem.attrib.get("y2"))

            line = slide.shapes.add_shape(
                MSO_SHAPE.LINE_INVERSE,
                px_to_emu(x1),
                px_to_emu(y1),
                px_to_emu(x2 - x1),
                px_to_emu(y2 - y1),
            )

            if "stroke" in elem.attrib:
                line.line.color.rgb = hex_to_rgb(elem.attrib["stroke"])

            if "marker-end" in elem.attrib:
                line.line.end_arrowhead = MSO_LINE_ARROWHEAD_STYLE.TRIANGLE

        # TEXT
        elif tag == "text":
            x = float(elem.attrib.get("x"))
            y = float(elem.attrib.get("y"))
            font_size = float(elem.attrib.get("font-size", 14))

            textbox = slide.shapes.add_textbox(
                px_to_emu(x),
                px_to_emu(y),
                px_to_emu(400),
                px_to_emu(100),
            )

            tf = textbox.text_frame
            tf.clear()

            full_text = "".join(elem.itertext())
            tf.text = full_text

            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    if "fill" in elem.attrib:
                        run.font.color.rgb = hex_to_rgb(elem.attrib["fill"])

    prs.save(output)
    print("Saved to", output)
