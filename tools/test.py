from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pptx.enum.shapes

# 检查是否存在，不存在则打印当前所有可用枚举名（方便调试）
try:
    from pptx.enum.shapes import MSO_LINE_ARROWHEAD_STYLE
except ImportError:
    import pptx.enum
    print("当前可用枚举:", dir(pptx.enum.shapes))